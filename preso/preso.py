from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Rebuild presentation (since previous one not accessible after reset)
prs = Presentation()

# Define slide layout (Title and Content)
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# --- Title Slide ---
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Banking Use Cases for Blockchain, AI & Agentic AI"
subtitle.text = "Integrating Blockchain, AI, and Agentic AI for Smarter Banking Solutions"

# --- Blockchain Use Cases ---
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Leading Blockchain Use Cases"
content.text = (
    "• Faster Payments: Blockchain enables instant payments (including cross-border)\n"
    "• Clearance and Settlement: Direct settlement via digital ledger\n"
    "• Trade Finance and Asset Tokenization: Streamlines complex trade finance\n"
    "• Regulatory Reporting and Compliance: Automates reporting & ensures data integrity\n"
    "• Fraud and Security: Immutable ledger supports anti-fraud strategies\n"
    "• Decentralized Lending & Borrowing: Peer-to-peer transparent lending\n"
    "• Central Bank Digital Currencies (CBDC): Blockchain-based digital money"
)
slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CLOUD, Inches(5.5), Inches(1.5), Inches(1.2), Inches(1.2)).fill.fore_color.rgb = RGBColor(0, 176, 240)

# --- AI Use Cases ---
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Leading AI Use Cases"
content.text = (
    "• Fraud Detection: Detect anomalies in transaction patterns\n"
    "• AI Chatbots & Customer Service: Virtual agents offer 24/7 support\n"
    "• Credit Risk & Lending Analytics: Assess creditworthiness using ML models\n"
    "• KYC & Compliance Automation: Accelerates checks, reduces manual work\n"
    "• Predictive Analytics (Markets): Tailored investment advice, forecast risks\n"
    "• Contract Intelligence: Automate document review and compliance checking"
)
slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT, Inches(5.5), Inches(1.5), Inches(1.2), Inches(1.2)).fill.fore_color.rgb = RGBColor(255, 192, 0)

# --- Integrated Blockchain + AI Use Cases ---
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Integrated Blockchain & AI Use Cases"
content.text = (
    "• AI-Driven Smart Contract Automation: Real-time triggers for agreements\n"
    "• Combined Fraud and Cybersecurity: AI scans blockchain data for anomalies\n"
    "• Enhanced AML: AI + Blockchain transparency to detect suspicious transactions\n"
    "• Financial Inclusion & Analytics: AI expands access to financial services"
)
slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(5.5), Inches(1.5), Inches(1.2), Inches(1.2)).fill.fore_color.rgb = RGBColor(146, 208, 80)

# --- Agentic AI Enhancement ---
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Agentic AI Across All Areas"
content.text = (
    "• Self-Driving Decision Making: Agentic AI automates end-to-end workflows\n"
    "• Dynamic Risk Management: Continuously adapts fraud, AML & compliance strategies\n"
    "• Autonomous Smart Contracts: Proactively executes and manages agreements\n"
    "• Customer-Centric Banking: Personalized experiences via proactive virtual agents\n"
    "• Regulatory & Market Adaptability: Real-time adjustments to policy & market changes\n"
    "• Financial Inclusion: Expands reach with intelligent, adaptive financial services"
)
slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, Inches(5.5), Inches(1.5), Inches(1.2), Inches(1.2)).fill.fore_color.rgb = RGBColor(255, 0, 0)

# --- Quick Table ---
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Quick Table: Top Banking Use Cases"
content.text = (
    "Blockchain Use Case → AI Use Case → Integrated Blockchain + AI → Agentic AI Enhancement\n\n"
    "• Instant cross-border payments → Fraud detection → AI-driven smart contracts → Autonomous settlement agents\n"
    "• Clearing & settlement → Customer service/chatbots → AML & anomaly detection → Self-adjusting compliance agents\n"
    "• Trade finance, asset tokenization → Credit risk assessment → Security & compliance analytics → Adaptive risk-mitigation agents"
)

# Save final version with visuals
output_path_visual = "/mnt/data/Banking_Blockchain_AI_AgenticAI_Visual.pptx"
prs.save(output_path_visual)

output_path_visual
